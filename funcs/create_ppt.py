# Creating powerpoint presentations using the python-pptx package

import collections
import collections.abc
from pptx import Presentation

X = Presentation()

Layout = X.slide_layouts[0]
first_slide = X.slides.add_slide(Layout)

first_slide.shapes.title.text = "Test creating a powerpoint using Python"
first_slide.placeholders[1].text = "Success!"

X.save("First_presentation.pptx")
